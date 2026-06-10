#!/usr/bin/env python3
"""
build_deck_4th.py — rebuild capstone_presentation.pptx for the 4th-cap story.

Fresh 13-slide deck in the established template (navy/teal/amber, Calibri, numbered eyebrows,
accent-bar cards, takeaway line). Native cards (not pasted charts) keep the dark theme clean.
Narrative: intro/papers/data (kept) → the collapse → categorize → honest ladder → NOW bake-off
→ free-geodata ceiling-breaker → forecast → verdict → contribution. All numbers from FINDINGS.md.

Run: ../../working/.venv/bin/python build_deck_4th.py
"""
import os
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
C = dict(navy="0D1B2A", card="122030", teal="14A38B", amber="E8A020", blue="5BA4CF",
         body="A0B8C4", muted="7A9BAD", light="E8F0EE", white="FFFFFF", good="46C98A", bad="EC5A4C")
def rgb(k): return RGBColor.from_string(C.get(k, k))
F = "Calibri"

prs = Presentation(); prs.slide_width = IN(13.333); prs.slide_height = IN(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb("navy")
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, IN(13.333), IN(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = rgb("teal"); bar.line.fill.background()
    foot = box(s, 0.0, 7.05, 13.333, 0.32, "GUC Masters   ·   Smart City Capstone   ·   2026", 10.5, "muted", align=PP_ALIGN.CENTER)
    return s

def box(s, x, y, w, h, text="", size=12, color="body", bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(text, str): text = [(text, color, bold)]
    for i, seg in enumerate(text):
        t, c, b = (seg + (False,))[:3] if isinstance(seg, tuple) else (seg, color, bold)
        r = p.add_run(); r.text = t; r.font.name = F; r.font.size = Pt(size); r.font.bold = b; r.font.color.rgb = rgb(c)
    return tb

def head(s, eyebrow, title):
    box(s, 0.60, 0.42, 12.0, 0.32, eyebrow, 12, "teal", True)
    box(s, 0.60, 0.82, 12.10, 1.0, title, 31, "white", True)

def rect(s, x, y, w, h, fill="card"):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    r.fill.solid(); r.fill.fore_color.rgb = rgb(fill); r.line.fill.background()
    r.shadow.inherit = False
    return r

def card(s, x, y, w, h, accent, label, body, lab_color=None, label_sz=14.5, body_sz=12):
    rect(s, x, y, w, h, "card"); rect(s, x, y, 0.07, h, accent)
    inner = []
    if label: inner.append((label + "   ", lab_color or accent, True))
    if body: inner.append((body, "body", False))
    box(s, x + 0.28, y + 0.10, w - 0.45, h - 0.18, inner, label_sz, anchor=MSO_ANCHOR.MIDDLE)

def takeaway(s, text, y=6.45):
    box(s, 0.60, y, 12.10, 0.45, [("→  ", "teal", True), (text, "muted", False)], 13.5, anchor=MSO_ANCHOR.MIDDLE)

def stat(s, x, y, w, val, lab, color="teal"):
    box(s, x, y, w, 0.55, val, 27, color, True, PP_ALIGN.CENTER)
    box(s, x, y + 0.62, w, 0.55, lab, 11, "body", False, PP_ALIGN.CENTER)

# ============================================================ SLIDES
# 1 — title
s = slide()
box(s, 0.60, 0.42, 12, 0.32, "SMART CITY CAPSTONE  ·  GUC MASTERS", 12, "teal", True)
box(s, 0.60, 1.55, 12.1, 1.4, [("Categorize-then-Predict\n", "white", True), ("Trustworthy PM2.5 where there is no sensor — and a few days ahead", "body", False)], 34, "white", True)
box(s, 0.60, 3.30, 12.1, 0.5, "Kampala, Uganda  ·  39 low-cost sensors  ·  free data only", 15, "amber", True)
card(s, 0.60, 4.20, 12.10, 1.7, "teal", "The question",
     "Can a categorize-then-predict pipeline give an urban planner trustworthy PM2.5 where there is no sensor, and a few days ahead, without collapsing? This deck reports the honest answer — out-of-site and out-of-time, negative results included.", body_sz=13)

# 2 — research journey (kept context)
s = slide(); head(s, "01  ·  RESEARCH JOURNEY", "From 60+ Papers Down to 4")
for i,(t,b,ac) in enumerate([
    ("60+ screened","Air-quality + ML + remote-sensing literature for low-resource cities.","blue"),
    ("4 anchors kept","Each builds on the last — precedent, levers, honest evaluation.","teal"),
    ("1 gap","Models look great on a random split, then collapse at a new place.","amber")]):
    card(s, 0.60, 2.35+i*1.05, 12.10, 0.9, ac, t, b)
takeaway(s, "The gap — generalization to unseen places — is exactly what the 4th cap attacks.")

# 3 — the 4 papers (kept)
s = slide(); head(s, "02  ·  THE FOUR PAPERS", "Each Paper Builds on the Last")
papers=[("Adong 2025","Kampala PM2.5 from satellite AOD alone — weather deliberately excluded. Our precedent.","teal"),
        ("Weather lever","ERA5 weather is the cheap, generalizable signal the precedent left out.","blue"),
        ("Honest evaluation","African spatial-CV ~0.13 vs random >0.90 — collapse is the norm, not the exception.","amber"),
        ("Land-use regression","Road density / distance-to-road / terrain transfer to new locations.","good")]
for i,(t,b,ac) in enumerate(papers):
    card(s, 0.60, 2.30+i*1.02, 12.10, 0.88, ac, t, b)
takeaway(s, "Together: cheap data can work — but only if evaluated honestly and given place-structure.")

# 4 — the reframe
s = slide(); head(s, "03  ·  THE REFRAME", "From “Which Lever?” to “Categorize-then-Predict”")
card(s, 0.60, 2.30, 5.85, 1.5, "muted", "Old framing (3rd cap)", "“Which cheap data lever helps estimate PM2.5?” → a finding (weather helps, LST is inert), but the tool collapses at a sensor-less place.", label_sz=14)
card(s, 6.75, 2.30, 5.95, 1.5, "teal", "New framing (4th cap)", "A clinic, not one doctor: a triage nurse sorts each place into a kind, then the specialist for that kind predicts.", label_sz=14)
box(s, 0.60, 4.05, 12.1, 0.35, "Two metrics — kept separate end to end:", 14, "light", True)
card(s, 0.60, 4.50, 5.85, 1.4, "blue", "NOW · nowcast", "PM2.5 today at a sensor-less place. Honest test = leave-sites-out. Bar = the blind pooled model.", label_sz=14)
card(s, 6.75, 4.50, 5.95, 1.4, "amber", "FUTURE · forecast", "PM2.5 +1…7 days at a known place. Honest test = chronological. Bar = persistence (“same as yesterday”).", label_sz=14)
takeaway(s, "Different question, different honest test, different bar, different best model.")

# 5 — the data
s = slide(); head(s, "04  ·  THE DATA", "39 Sensors · One Merged Table · Free Inputs Only")
for i,(v,l,c) in enumerate([("39","AirQo sensors","teal"),("13,944","site-days","blue"),("~84%","daily coverage","amber"),("38","µg/m³ mean PM2.5","good")]):
    stat(s, 0.70+i*3.05, 2.35, 2.8, v, l, c)
card(s, 0.60, 4.05, 12.10, 1.05, "teal", "Predictors (free, available anywhere)", "ERA5 weather (temp · dew-point · wind · pressure · rain) + location + calendar. LST carried but proven inert.", label_sz=13.5)
card(s, 0.60, 5.20, 12.10, 1.0, "bad", "Banned as predictors", "pm10 (same instrument as the target, corr ≈ 0.9) and n_obs (a by-product of the daily mean) — neither exists at a sensor-less place, so using them would manufacture a fake win.", lab_color="bad", label_sz=13.5)
takeaway(s, "Honest features first: nothing a planner could not actually obtain at a new address.")

# 6 — the collapse
s = slide(); head(s, "05  ·  THE COLLAPSE", "Looks Great — Until You Point It Somewhere New")
for i,(v,l,c) in enumerate([("0.81","leaky random-split R²","muted"),("0.31","honest new-site R²","bad"),("~0","one new-site fold","bad")]):
    stat(s, 0.90+i*4.0, 2.35, 3.6, v, l, c)
card(s, 0.60, 4.05, 12.10, 1.35, "amber", "Why it matters",
     "The headline R²≈0.81 is a leaky random split — it tests on the same sensors it trained on (gap-filling). Graded honestly (hide whole sensors, or forecast the real future) it collapses, and a “same-as-yesterday” guess beats it out-of-time. A planner needs PM2.5 where there is no sensor.", body_sz=12.5)
takeaway(s, "The whole 4th cap exists to stop this collapse — without faking the test.")

# 7 — categorize
s = slide(); head(s, "06  ·  CATEGORIZE (CP-A)", "39 Sensors → 3 Behavioural Kinds of Place")
rows=[("24","low-pollution / well-ventilated — the city baseline","blue","PM2.5 ≈ 31.7"),
      ("12","high-pollution / well-ventilated — hotter spots, still windy","amber","PM2.5 ≈ 47.5"),
      ("3","moderate / stagnant — low wind + humid, a robust micro-regime","teal","PM2.5 ≈ 46.5")]
for i,(n,d,ac,m) in enumerate(rows):
    y=2.35+i*0.95
    rect(s,0.60,y,12.10,0.8,"card"); rect(s,0.60,y,0.07,0.8,ac)
    box(s,0.85,y,1.1,0.8,n,26,ac,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    box(s,2.05,y,8.2,0.8,d,13,"body",False,anchor=MSO_ANCHOR.MIDDLE)
    box(s,10.4,y,2.1,0.8,m,13.5,ac,True,PP_ALIGN.RIGHT,MSO_ANCHOR.MIDDLE)
card(s, 0.60, 5.30, 12.10, 0.95, "good", "Robust + portable", "Same 3 stagnant sites at every k; k-means agrees with Ward. Categories are spatially intermixed — they capture behaviour, not geography. A new site is categorized from its signature, snapped to the nearest centroid.", label_sz=13)
takeaway(s, "Kampala is homogeneous in level; the structure is low-vs-high pollution + 3 stagnant outliers.")

# 8 — the honest ladder
s = slide(); head(s, "07  ·  THE HONEST LADDER (CP-B)", "Knowing the Category Beats the Blind Model")
box(s,0.60,2.20,12.1,0.32,"New-site R² (leave-sites-out) — climb only if the rung below holds:",13.5,"light",True)
bars=[("blind (no category)","0.316","bad",5.6),("+ category · k=4 (tuned)","0.558","teal",9.4),("+ category + geodata · k=4","0.564","good",9.7)]
for i,(t,v,ac,w) in enumerate(bars):
    y=2.70+i*0.78
    rect(s,0.60,y,w,0.62,"card"); rect(s,0.60,y,0.07,0.62,ac)
    box(s,0.85,y,w-1.6,0.62,t,13.5,"light",True,anchor=MSO_ANCHOR.MIDDLE)
    box(s,0.60+w-1.5,y,1.4,0.62,v,17,ac,True,PP_ALIGN.RIGHT,MSO_ANCHOR.MIDDLE)
card(s, 0.60, 5.15, 5.85, 1.1, "teal", "Significant", "Tuning categories to k=4 (by honest R², not silhouette) lifts the deploy-only nowcast 0.48→0.56. Paired bootstrap +1.07 µg/m³, CI [+0.33, +1.90].", label_sz=13)
card(s, 6.75, 5.15, 5.95, 1.1, "amber", "Soft routing wins", "Category-as-a-feature tolerates a wrong category; mixture-of-experts hard-routes and collapses (0.50→0.30) under realistic cold-start.", label_sz=13)
takeaway(s, "Ceiling: leave a whole category out → negative. Interpolates among known kinds, can’t extrapolate.")

# 9 — NOW bake-off
s = slide(); head(s, "08  ·  NOW BAKE-OFF (CP-D)", "The Model Is Not the Bottleneck")
box(s,0.60,2.20,12.1,0.32,"New-site R² · same +category features · four learners:",13.5,"light",True)
for i,(m,v,c) in enumerate([("RandomForest","0.529","good"),("XGBoost","0.482","teal"),("LightGBM","0.469","blue"),("CatBoost","0.465","amber")]):
    stat(s, 0.70+i*3.05, 2.70, 2.8, v, m, c)
card(s, 0.60, 4.30, 5.85, 1.05, "good", "Best: RandomForest", "Lowest leaky score yet highest new-site score — it overfits least, transfers best. But all 4 cluster (0.47–0.53).", label_sz=13)
card(s, 6.75, 4.30, 5.95, 1.05, "good", "Conformal — fixed", "Calibrating on held-out SITES (group conformal) lifts new-site 90%-interval coverage 74.8% → 85.8%. Bands widen honestly.", label_sz=13)
takeaway(s, "The category feature (+0.16) moved R² far more than any model swap. Features > learner.")

# 10 — break the ceiling
s = slide(); head(s, "09  ·  BREAK THE CEILING", "Free Map Data Nowcasts a Sensor-less Site")
for i,(v,l,c) in enumerate([("0.32","blind","muted"),("0.42","+ geodata · ZERO sensors","blue"),("0.55","+ category + geodata","good")]):
    stat(s, 0.90+i*4.0, 2.35, 3.6, v, l, c)
card(s, 0.60, 4.05, 12.10, 1.15, "blue", "Land-use regression, sourced free",
     "OSM road density + distance-to-nearest-major-road (Overpass, 43,870 ways) + SRTM elevation. Static, need ZERO pollution observations. dist-to-major-road is the clearest signal (closer → higher PM2.5).", body_sz=12.5)
card(s, 0.60, 5.30, 12.10, 0.95, "amber", "A complement, not a substitute", "Geodata can’t name the category (54% < 62% baseline) — it adds independent signal. Combined is best, and most stable (±0.046).", label_sz=13)
takeaway(s, "The bridge from “a finding” to “a tool a planner can point at any address.”")

# 11 — forecast
s = slide(); head(s, "10  ·  FORECAST (CP-D)", "A Real Forecaster Holds R²≈0.43 a Week Out")
box(s,0.60,2.20,12.1,0.32,"R² by horizon (chronological, out-of-time):",13.5,"light",True)
hdr=["horizon","persistence","GBM-on-lags","LSTM"]
data=[["+1 day","0.34","0.43","0.51"],["+3 days","0.11","0.37","0.48"],["+7 days","0.02","0.33","0.43"]]
cw=[2.6,3.1,3.1,3.0]; x0=0.60; y0=2.70
for j,htx in enumerate(hdr):
    box(s,x0+sum(cw[:j]),y0,cw[j],0.4,htx,13,"muted",True,PP_ALIGN.CENTER)
for i,row in enumerate(data):
    yy=y0+0.5+i*0.62
    rect(s,0.60,yy,sum(cw),0.54,"card")
    for j,val in enumerate(row):
        col="light" if j==0 else ("bad" if j==1 else ("amber" if j==2 else "good"))
        box(s,x0+sum(cw[:j]),yy,cw[j],0.54,val,14 if j else 13,col,j!=0,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
card(s, 0.60, 5.35, 12.10, 0.95, "amber", "LSTM earns its weight", "Global LSTM + category embedding beats the GBMs at every horizon (+0.04…+0.12), lead widening with horizon. All beat persistence. Correction: forecasting does NOT decay to ~0 by +7d.", label_sz=13)
takeaway(s, "Persistence-beats-model was true only for the nowcast model — not a purpose-built forecaster.")

# 12 — safety tiers (NEW SCOPE — the planner's decision)
s = slide(); head(s, "11  ·  SAFETY TIERS", "From a Number to a Safety Status")
box(s,0.60,2.05,12.1,0.32,"EPA PM2.5 safety bands — what a city official actually needs:",13.5,"light",True)
for i,(v,l,c) in enumerate([("Elevated","≤ 35.4 µg/m³","blue"),("High","35.5 – 55.4","amber"),("Dangerous","> 55.4","bad")]):
    stat(s, 0.90+i*4.0, 2.45, 3.6, v, l, c)
box(s,0.60,3.95,12.1,0.32,"Graded by the BAND it lands in — same honest tests:",13.5,"light",True)
card(s,0.60,4.35,5.85,1.05,"good","New-site (no sensor)","Exact tier 68%  ·  within one tier 98.5%  ·  vs 52% majority baseline. A 2-tier error (a Dangerous area called safe) for only 3.6% of days.",label_sz=13)
card(s,6.75,4.35,5.95,1.05,"teal","Forecast +1→+7d","Exact tier 70% → 68%  ·  within one tier ~98%  ·  beats the persistence-band baseline at every horizon.",label_sz=13)
takeaway(s,"R² punished 47-vs-52 misses; in safety tiers those are the SAME band — so the tool is deployable.")

# 13 — danger alerting (NEW SCOPE)
s = slide(); head(s, "12  ·  DANGER ALERTING", "Catching the Days That Matter — a Planner's Dial")
box(s,0.60,2.05,12.1,0.32,"Lower the Dangerous alert cutoff → catch more dangerous days, at a chosen false-alarm cost (honest new-site):",12,"light",True)
hdr=["alert if predicted >","dangerous-day recall","false-alarm rate"]
data=[["55.4 (strict EPA)","0.51","4%"],["50","0.67","8%"],["45   (suggested)","0.80","17%"],["40","0.91","30%"]]
cw=[4.3,4.0,4.0]; x0=0.60; y0=2.55
for j,htx in enumerate(hdr): box(s,x0+sum(cw[:j]),y0,cw[j],0.4,htx,12.5,"muted",True,PP_ALIGN.CENTER)
for i,row in enumerate(data):
    yy=y0+0.5+i*0.55; rect(s,0.60,yy,sum(cw),0.48,"card")
    for j,val in enumerate(row):
        col="light" if j==0 else ("good" if j==1 else "amber")
        box(s,x0+sum(cw[:j]),yy,cw[j],0.48,val,13,col,j!=0,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
card(s,0.60,5.35,6.0,0.98,"teal","Why it works","Dangerous days are a humid, stagnant weather regime (humidity + dew-point separate them most) — the model keys on weather to flag them.",label_sz=12.5)
card(s,6.75,5.35,5.95,0.98,"bad","Honest negative","A dedicated class-weighted detector did NOT beat simply lowering the regressor's threshold. The simpler lever wins.",lab_color="bad",label_sz=12.5)
takeaway(s,"The planner owns the operating point: 80% of dangerous days caught at a 45 µg/m³ alert cutoff.")

# 14 — verdict
s = slide(); head(s, "13  ·  THE VERDICT", "The Reframe Works — With Honest Limits")
card(s, 0.60, 2.30, 3.86, 2.2, "good", "What worked", "Category lifts new-site R² 0.32→0.48 (significant). Free geodata: 0.42 zero-deploy, 0.55 combined. LSTM holds 0.43 at +7d.", label_sz=13.5)
card(s, 4.66, 2.30, 3.86, 2.2, "bad", "What didn’t", "Mixture-of-experts collapses under cold-start. Category×month prior hurts under noise. Geodata can’t name the category.", lab_color="bad", label_sz=13.5)
card(s, 8.72, 2.30, 3.98, 2.2, "amber", "Bounded", "Leave-a-whole-category-out goes negative — can interpolate among known kinds, can’t extrapolate to an unseen regime.", label_sz=13.5)
card(s, 0.60, 4.70, 12.10, 1.4, "teal", "Two prior claims corrected (verified)", "1) Forecasting does not decay to ~0 by +7d — an artefact of using a nowcast model to forecast. 2) “Categorization can’t rescue a homogeneous city” is too pessimistic — it lifts new-site R² by +0.16; true only for unseen regimes.", body_sz=12.5)
takeaway(s, "Honesty over headlines — every win holds out-of-site and out-of-time, or it isn’t reported.")

# 15 — contribution / next
s = slide(); head(s, "14  ·  CONTRIBUTION & NEXT", "A Cheap, Honest PM2.5 Safety Tool for a Low-Resource City")
card(s, 0.60, 2.30, 12.10, 1.25, "teal", "The contribution", "Free data + behavioural categories + land-use geodata turn a collapsing model into a deployable safety-triage tool: classify a never-instrumented area's safety tier (68% exact / 98.5% within one tier) and forecast a week ahead — evaluated the way a planner will actually use it.", body_sz=12.5)
nexts=[("Done · safety tiers","band evaluation + planner-tunable danger alerting (80% recall)","good"),
       ("Done · sharper model","k=4 categories (0.56) + group conformal (86% coverage)","good"),
       ("Next · more places","39 sites is small — more sensors/categories sharpen routing","teal"),
       ("Next · push recall","cost curves to 90% danger recall · tighter conformal to 90%","amber")]
for i,(t,b,ac) in enumerate(nexts):
    x=0.60+(i%2)*6.10; y=3.75+(i//2)*1.15
    card(s, x, y, 5.95, 1.0, ac, t, b, label_sz=13)
takeaway(s, "Reproducible: fixed seeds · honest splits · scripts + result CSVs in 4th-cap-handover/analysis/.")

import sys
out = os.path.join(HERE, "..", "docs", os.environ.get("DECK_OUT", "capstone_presentation.pptx"))
prs.save(out)
print(f"saved {out}  ·  {len(prs.slides)} slides")
